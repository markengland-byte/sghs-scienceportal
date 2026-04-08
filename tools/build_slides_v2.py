"""
Slide Deck Builder v2 — PDF pages as slide backgrounds with click-to-reveal answers.

Each PDF page becomes one slide. The page image is the background.
Answer textboxes are positioned over the blanks and appear on click.

Usage:
    python tools/build_slides_v2.py
"""

import sys
import io
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace')

TEAL = RGBColor(0x0B, 0x8F, 0x8C)
NAVY = RGBColor(0x0F, 0x22, 0x40)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREEN = RGBColor(0x16, 0xA3, 0x4A)
RED = RGBColor(0xDC, 0x26, 0x26)

# Slide dimensions (standard 4:3 works best for PDF pages)
SLIDE_W = 10.0
SLIDE_H = 7.5

# PDF page dimensions in pixels (at 2x render)
PDF_PX_W = 1224
PDF_PX_H = 1584

# The page image is scaled to fill the slide width, maintaining aspect ratio
# Image height on slide = SLIDE_W * (PDF_PX_H / PDF_PX_W)
IMG_H_ON_SLIDE = SLIDE_W * (PDF_PX_H / PDF_PX_W)  # ~12.94" — taller than slide
# So we scale to fit height instead: SLIDE_H * (PDF_PX_W / PDF_PX_H) for width
IMG_W_ON_SLIDE = SLIDE_H * (PDF_PX_W / PDF_PX_H)  # ~5.80"
# Center horizontally
IMG_LEFT = (SLIDE_W - IMG_W_ON_SLIDE) / 2


def px_to_inches(px_x, px_y):
    """Convert PDF pixel coordinates (at 2x) to slide inches."""
    x = IMG_LEFT + (px_x / PDF_PX_W) * IMG_W_ON_SLIDE
    y = (px_y / PDF_PX_H) * SLIDE_H
    return x, y


def add_appear_animations(slide, shape_ids):
    """Add click-to-appear animations for shapes. Each click reveals the next."""
    if not shape_ids:
        return

    timing = slide.element.makeelement(qn('p:timing'), {})
    tn_lst = timing.makeelement(qn('p:tnLst'), {})
    par = tn_lst.makeelement(qn('p:par'), {})
    c_tn = par.makeelement(qn('p:cTn'), {'id': '1', 'dur': 'indefinite', 'restart': 'never', 'nodeType': 'tmRoot'})
    child_lst = c_tn.makeelement(qn('p:childTnLst'), {})

    seq = child_lst.makeelement(qn('p:seq'), {'concurrent': '1', 'nextAc': 'seek'})
    seq_ctn = seq.makeelement(qn('p:cTn'), {'id': '2', 'dur': 'indefinite', 'nodeType': 'mainSeq'})
    seq_child = seq_ctn.makeelement(qn('p:childTnLst'), {})

    anim_id = 3
    for sp_id in shape_ids:
        par2 = seq_child.makeelement(qn('p:par'), {})
        par2_ctn = par2.makeelement(qn('p:cTn'), {'id': str(anim_id), 'fill': 'hold'})
        anim_id += 1
        st = par2_ctn.makeelement(qn('p:stCondLst'), {})
        st.append(st.makeelement(qn('p:cond'), {'delay': '0'}))
        par2_ctn.append(st)

        par2_child = par2_ctn.makeelement(qn('p:childTnLst'), {})
        par3 = par2_child.makeelement(qn('p:par'), {})
        par3_ctn = par3.makeelement(qn('p:cTn'), {'id': str(anim_id), 'fill': 'hold'})
        anim_id += 1
        st2 = par3_ctn.makeelement(qn('p:stCondLst'), {})
        st2.append(st2.makeelement(qn('p:cond'), {'delay': '0'}))
        par3_ctn.append(st2)

        par3_child = par3_ctn.makeelement(qn('p:childTnLst'), {})
        set_node = par3_child.makeelement(qn('p:set'), {})
        set_cbn = set_node.makeelement(qn('p:cBhvr'), {})
        set_ctn_inner = set_cbn.makeelement(qn('p:cTn'), {'id': str(anim_id), 'dur': '1', 'fill': 'hold'})
        anim_id += 1
        st3 = set_ctn_inner.makeelement(qn('p:stCondLst'), {})
        st3.append(st3.makeelement(qn('p:cond'), {'delay': '0'}))
        set_ctn_inner.append(st3)
        set_cbn.append(set_ctn_inner)

        tgt_el = set_cbn.makeelement(qn('p:tgtEl'), {})
        tgt_el.append(tgt_el.makeelement(qn('p:spTgt'), {'spid': str(sp_id)}))
        set_cbn.append(tgt_el)

        to_node = set_node.makeelement(qn('p:to'), {})
        to_node.append(to_node.makeelement(qn('p:strVal'), {'val': 'visible'}))
        set_node.append(set_cbn)
        set_node.append(to_node)

        par3_child.append(set_node)
        par3_ctn.append(par3_child)
        par3.append(par3_ctn)
        par2_child.append(par3)
        par2_ctn.append(par2_child)
        par2.append(par2_ctn)
        seq_child.append(par2)

    seq_ctn.append(seq_child)
    for evt_name, lst_name in [('onPrev', 'p:prevCondLst'), ('onNext', 'p:nextCondLst')]:
        cond_lst = seq.makeelement(qn(lst_name), {})
        cond = cond_lst.makeelement(qn('p:cond'), {'evt': evt_name, 'delay': '0'})
        tgt = cond.makeelement(qn('p:tgtEl'), {})
        tgt.append(tgt.makeelement(qn('p:sldTgt'), {}))
        cond.append(tgt)
        cond_lst.append(cond)
        seq.append(cond_lst)

    seq.insert(0, seq_ctn)
    child_lst.append(seq)
    c_tn.append(child_lst)
    par.append(c_tn)
    tn_lst.append(par)
    timing.append(tn_lst)
    slide.element.append(timing)


def add_answer_box(slide, x_in, y_in, w_in, h_in, text, font_size=11, color=TEAL, bg_white=True):
    """Add a rounded answer box positioned over a blank. Returns shape for animation."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(x_in), Inches(y_in), Inches(w_in), Inches(h_in)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE if bg_white else RGBColor(0xE0, 0xF6, 0xF6)
    shape.line.color.rgb = TEAL
    shape.line.width = Pt(1.5)
    shape.text_frame.word_wrap = True
    shape.text_frame.margin_left = Inches(0.05)
    shape.text_frame.margin_right = Inches(0.05)
    shape.text_frame.margin_top = Inches(0.02)
    shape.text_frame.margin_bottom = Inches(0.02)
    p = shape.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = True
    p.font.color.rgb = color
    p.font.name = 'Calibri'
    p.alignment = PP_ALIGN.CENTER
    return shape


def build_page_slide(prs, page_num, page_img_path, answers, notes=''):
    """Build a slide with the PDF page as background and answer overlays.

    answers: list of (x_px, y_px, w_px, h_px, answer_text, font_size)
        Coordinates are in PDF pixels at 2x resolution.
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    # Set white background
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = WHITE

    # Add page image centered
    slide.shapes.add_picture(
        str(page_img_path),
        Inches(IMG_LEFT), Inches(0),
        width=Inches(IMG_W_ON_SLIDE), height=Inches(SLIDE_H)
    )

    # Add answer overlays
    answer_ids = []
    for ans in answers:
        x_px, y_px, w_px, h_px, text = ans[:5]
        font_size = ans[5] if len(ans) > 5 else 11

        # Convert pixel coords to slide inches
        x_in, y_in = px_to_inches(x_px, y_px)
        w_in = (w_px / PDF_PX_W) * IMG_W_ON_SLIDE
        h_in = (h_px / PDF_PX_H) * SLIDE_H

        shape = add_answer_box(slide, x_in, y_in, w_in, h_in, text, font_size)
        answer_ids.append(shape.shape_id)

    if answer_ids:
        add_appear_animations(slide, answer_ids)

    if notes:
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = notes

    return slide


def build_ch1():
    import importlib.util
    spec = importlib.util.spec_from_file_location('ch1_answers', Path(__file__).parent / 'ch1_answers.py')
    mod = importlib.util.module_from_spec(spec)
    spec.loader.exec_module(mod)
    ANSWERS = mod.ANSWERS

    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)

    pages_dir = Path(__file__).parent.parent / 'slides' / 'images' / 'ch1' / 'pages'

    for page_num in range(1, 16):
        page_img = pages_dir / f'page-{page_num:02d}.png'
        if not page_img.exists():
            continue
        answers = ANSWERS.get(page_num, [])
        build_page_slide(prs, page_num, page_img, answers,
            notes=f'Page {page_num}: {len(answers)} blanks. Click to reveal each answer.')

    return prs


def _OLD(): pass; return  # noqa
_REMOVED = """
        # Suffix "___" means "the study of"
        (596, 152, 90, 28, 'ology', 11),
        # most life on Earth is too ___ for our eyes
        (710, 186, 100, 28, 'small', 11),
        # ___: the smallest, most basic unit of life
        (78, 214, 80, 28, 'Cell', 11),
        # ___: any individual form of life
        (78, 240, 130, 28, 'Organism', 11),
        # "___" (Life) in the bio/ology diagram
        (118, 290, 65, 28, 'Bio', 12),
        # (The ___ of)
        (370, 308, 70, 24, 'study', 10),
        # uni or multi cellular
        (310, 500, 55, 26, 'uni', 10),
        (420, 500, 85, 26, 'multi', 10),
        # consist of a ___ cell
        (430, 528, 80, 26, 'single', 10),
        # consist of ___ cells
        (365, 550, 150, 26, 'many / multiple', 10),
        # ___ cellular Organism (left diagram)
        (190, 588, 55, 24, 'Uni', 10),
        # ___ cellular Organism (right diagram)
        (468, 588, 75, 24, 'Multi', 10),
        # All living organisms are ___
        (360, 766, 150, 26, 'either uni- or multicellular', 9),
    ], notes='Page 1: Introduction to Biology. Click to reveal each blank answer.')

    # ═══════════════════════════════════════════════
    # PAGE 2: Characteristics of Life
    # ═══════════════════════════════════════════════
    build_page_slide(prs, 2, pages_dir / 'page-02.png', [
        # share ___ characteristics
        (278, 92, 55, 26, 'eight', 11),
        # 1) Composed of ___
        (188, 134, 70, 24, 'cells', 10),
        # 2) ___: use smaller structures
        (90, 176, 110, 24, 'Organization', 9),
        # 3) Respond to environmental ___
        (338, 218, 90, 24, 'stimuli', 10),
        # 4) Maintain ___: stabilize suitable
        (168, 258, 130, 24, 'homeostasis', 10),
        # 5) ___: capacity to produce life
        (580, 134, 110, 24, 'Reproduction', 9),
        # sexually or ___
        (740, 150, 90, 24, 'asexually', 10),
        # 6) Acquires & utilizes ___ from surroundings
        (722, 176, 80, 24, 'energy', 10),
        # 7) ___ information: DNA
        (585, 218, 80, 24, 'Genetic', 10),
        # 8) ___: changes in ___ over time
        (580, 258, 90, 24, 'Evolution', 10),
        (760, 268, 55, 22, 'DNA', 10),
        # NOTE: ___ are NOT considered alive
        (118, 330, 80, 28, 'Viruses', 11),
    ], notes='Page 2: 8 Characteristics of Life. Viruses are NOT alive.')

    # ═══════════════════════════════════════════════
    # PAGE 3: Life's Organizational Hierarchy
    # ═══════════════════════════════════════════════
    build_page_slide(prs, 3, pages_dir / 'page-03.png', [
        # composed of ___; largest scale is ___
        (390, 120, 70, 26, 'atoms', 11),
        (740, 120, 110, 26, 'biosphere', 10),
        # Molecules (blank label)
        (82, 216, 100, 24, 'Molecules', 10),
        # Organelles: Specialized ___ within cells
        (500, 248, 105, 22, 'structures', 9),
        # ___ (Cells label)
        (82, 280, 80, 24, 'Cells', 10),
        # Tissues: Group of ___ performing
        (438, 312, 70, 22, 'cells', 9),
        # ___ (Organs label)
        (82, 344, 80, 24, 'Organs', 10),
        # Group of ___ that perform
        (440, 346, 70, 22, 'tissues', 9),
        # Group of ___ working together
        (440, 380, 70, 22, 'organs', 9),
        # ___ (Organism label)
        (82, 414, 100, 24, 'Organism', 10),
        # All organisms of the ___ species
        (468, 446, 80, 22, 'same', 9),
        # ___ (Community label)
        (82, 478, 100, 24, 'Community', 10),
        # Multiple populations of ___ species
        (425, 480, 100, 22, 'different', 9),
        # Living community & the ___ surroundings
        (520, 514, 90, 22, 'nonliving', 9),
        # ___ (Biosphere label)
        (82, 546, 100, 24, 'Biosphere', 10),
    ], notes='Page 3: Hierarchy of Life pyramid. Small to large: atom → biosphere.')

    # ═══════════════════════════════════════════════
    # PAGE 4: Hierarchy continued + Emergent Properties
    # ═══════════════════════════════════════════════
    build_page_slide(prs, 4, pages_dir / 'page-04.png', [
        # Emergent properties: properties that ___ upon combining
        (482, 584, 80, 26, 'arise', 11),
        # the whole is ___ than the sum
        (500, 616, 80, 26, 'greater', 11),
        # ___ Properties (label in diagram)
        (820, 680, 90, 24, 'Emergent', 10),
    ], notes='Page 4: Hierarchy order practice + Emergent Properties.')

    # ═══════════════════════════════════════════════
    # PAGE 5: Emergent Properties practice (no blanks)
    # ═══════════════════════════════════════════════
    build_page_slide(prs, 5, pages_dir / 'page-05.png', [],
        notes='Page 5: Practice questions on emergent properties. No blanks.')

    # ═══════════════════════════════════════════════
    # PAGE 6: Natural Selection & Evolution
    # ═══════════════════════════════════════════════
    build_page_slide(prs, 6, pages_dir / 'page-06.png', [
        # well suited due to ___
        (530, 92, 110, 26, 'adaptation', 11),
        # enables organisms to ___ survival
        (400, 118, 100, 24, 'improve', 10),
        # ___: ability to ___ & ___
        (70, 142, 80, 24, 'Fitness', 10),
        (340, 142, 80, 24, 'survive', 10),
        (460, 142, 100, 24, 'reproduce', 10),
        # Natural ___
        (544, 332, 110, 28, 'Selection', 12),
        # environment ___ for organisms that are more "___"
        (440, 368, 80, 24, 'selects', 10),
        (720, 368, 45, 24, 'fit', 10),
        # Survival of the ___
        (368, 392, 80, 24, 'fittest', 10),
        # Natural Selection has ___ requirements
        (302, 418, 40, 24, 'two', 10),
        # 1) Genetic ___
        (188, 442, 80, 24, 'variation', 10),
        # 2) Selective ___
        (500, 442, 80, 24, 'pressure', 10),
        # Genetic ___ in giraffes
        (108, 658, 80, 22, 'variation', 9),
    ], notes='Page 6: Adaptation, fitness, natural selection, two requirements.')

    # ═══════════════════════════════════════════════
    # PAGE 7: Evolution
    # ═══════════════════════════════════════════════
    build_page_slide(prs, 7, pages_dir / 'page-07.png', [
        # Evolution: changes in the ___ of a population
        (308, 272, 55, 26, 'DNA', 11),
        # can occur in a ___ of ways
        (368, 302, 80, 24, 'number', 10),
        # responsible for life's ___
        (528, 326, 80, 24, 'diversity', 10),
        # Mostly ___ crickets (left)
        (118, 500, 60, 22, 'light', 9),
        # Mostly ___ crickets (right)
        (688, 500, 60, 22, 'dark', 9),
    ], notes='Page 7: Evolution definition + cricket example.')

    # ═══════════════════════════════════════════════
    # PAGE 8: Introduction to Taxonomy
    # ═══════════════════════════════════════════════
    build_page_slide(prs, 8, pages_dir / 'page-08.png', [
        # ___: the branch of science that classifies
        (62, 92, 110, 28, 'Taxonomy', 11),
        # ___ categories
        (78, 120, 50, 24, 'Eight', 10),
        # ___ Inclusive (left - Most)
        (72, 162, 40, 20, 'Most', 9),
        # ___ Inclusive (right - Least)
        (830, 162, 40, 20, 'Least', 9),
        # D___ear K___ing P___hilip etc (mnemonic)
        (108, 236, 30, 18, 'D', 9),
        (220, 236, 30, 18, 'K', 9),
        (332, 236, 30, 18, 'P', 9),
        (560, 236, 30, 18, 'O', 9),
        (665, 236, 30, 18, 'F', 9),
        (770, 236, 30, 18, 'G', 9),
        (870, 236, 30, 18, 'S', 9),
        # 1) Bacteria  2) Archaea  3) Eukarya
        (72, 380, 100, 26, 'Bacteria', 10),
        (72, 406, 100, 26, 'Archaea', 10),
        (72, 444, 100, 26, 'Eukarya', 10),
        # Consist of ___ cells (___ a nucleus) — prokaryotes
        (360, 388, 100, 22, 'prokaryotic', 9),
        (520, 388, 60, 22, 'lack', 9),
        # Consist of ___ cells (contain a nucleus) — eukaryotes
        (360, 444, 100, 22, 'eukaryotic', 9),
        # ___-cellular organisms (prokaryotes)
        (812, 430, 40, 20, 'Uni', 9),
        # ___-cellular or ___-cellular (eukaryotes)
        (790, 486, 40, 20, 'Uni', 8),
        (810, 510, 55, 20, 'Multi', 8),
    ], notes='Page 8: Taxonomy, 8 categories, 3 domains, pro/eukaryote.')

    # ═══════════════════════════════════════════════
    # PAGE 9: Taxonomy continued — Kingdoms
    # ═══════════════════════════════════════════════
    build_page_slide(prs, 9, pages_dir / 'page-09.png', [
        # subdivided into ___
        (462, 276, 90, 24, 'kingdoms', 10),
        # Domain Eukarya has ___ kingdoms
        (372, 300, 40, 24, 'four', 10),
        # Kingdom ___ (Fungi)
        (600, 332, 60, 22, 'Fungi', 10),
        # ___ cellular (Animalia)
        (120, 426, 50, 20, 'Multi', 9),
        # ___ (Protista) — Unicellular or Multicellular
        (760, 388, 80, 20, 'Unicellular', 8),
    ], notes='Page 9: 4 kingdoms of Eukarya + tree of life diagram.')

    # ═══════════════════════════════════════════════
    # PAGE 10: Energy Acquisition
    # ═══════════════════════════════════════════════
    build_page_slide(prs, 10, pages_dir / 'page-10.png', [
        # ___ classes based on energy
        (372, 92, 48, 24, 'three', 10),
        # 1) ___ (___trophs)
        (72, 118, 100, 24, 'Producers', 10),
        (202, 118, 50, 24, 'Auto', 10),
        # 2) ___ (___trophs)
        (72, 144, 100, 24, 'Consumers', 10),
        (202, 144, 60, 24, 'Hetero', 10),
        # 3) ___: acquire energy from wastes
        (72, 170, 110, 24, 'Decomposers', 10),
        # Most energy from the ___
        (430, 200, 50, 24, 'sun', 10),
        # some energy is lost as ___
        (468, 222, 50, 24, 'heat', 10),
        # ___ (Autotroph label in diagram)
        (348, 558, 85, 22, 'Autotroph', 9),
        # ___ (Heterotroph label)
        (612, 558, 95, 22, 'Heterotroph', 9),
        # ___ (Decomposer label)
        (478, 688, 95, 22, 'Decomposer', 9),
    ], notes='Page 10: Autotrophs, heterotrophs, decomposers. Energy from sun, lost as heat.')

    # ═══════════════════════════════════════════════
    # PAGE 11: Scientific Method
    # ═══════════════════════════════════════════════
    build_page_slide(prs, 11, pages_dir / 'page-11.png', [
        # subject to the ___ method
        (578, 118, 100, 26, 'scientific', 10),
        # procedure used to ___ questions, ___ ideas, & ___ scientific knowledge
        (292, 144, 100, 24, 'investigate', 9),
        (418, 144, 55, 24, 'test', 9),
        (510, 144, 55, 24, 'build', 9),
        # starts with an ___ & a ___
        (340, 170, 100, 24, 'observation', 9),
        (500, 170, 80, 24, 'question', 9),
        # 5) Collect & Interpret ___
        (358, 368, 50, 22, 'Data', 9),
        # 3) Formulate ___ & make Prediction
        (668, 358, 80, 22, 'Hypothesis', 8),
        # 4) Design & Conduct ___
        (470, 430, 90, 22, 'Experiment', 8),
        # ___: an expected outcome
        (62, 520, 90, 24, 'Prediction', 10),
        # Predictions only address "___"
        (410, 544, 50, 22, 'what', 9),
        # ___: proposed & testable explanation
        (62, 580, 90, 24, 'Hypothesis', 10),
        # "___" it will happen?
        (530, 604, 45, 22, 'why', 9),
        # hypothesis ___ a prediction
        (322, 628, 80, 22, 'includes', 9),
        # ___: a testable & ___ hypothesis
        (62, 660, 75, 24, 'Theory', 10),
        (268, 660, 65, 22, 'broad', 9),
        # ___: (Prediction label in motorcycle diagram)
        (292, 706, 80, 18, 'Prediction', 8),
        # ___: (Hypothesis label)
        (528, 706, 80, 18, 'Hypothesis', 8),
        # ___: (Theory label)
        (780, 706, 55, 18, 'Theory', 8),
    ], notes='Page 11: Scientific method 7 steps, prediction vs hypothesis vs theory.')

    # ═══════════════════════════════════════════════
    # PAGE 12: Scientific Method practice + 3 Theories
    # ═══════════════════════════════════════════════
    build_page_slide(prs, 12, pages_dir / 'page-12.png', [
        # Label: prediction, hypothesis, theory
        (62, 392, 80, 24, 'Prediction', 10),
        (62, 420, 80, 24, 'Hypothesis', 10),
        (62, 448, 65, 24, 'Theory', 10),
        # There are ___ basic theories
        (178, 500, 48, 24, 'three', 10),
        # 1) ___
        (148, 544, 80, 22, 'Cell Theory', 9),
        # 2) ___
        (148, 580, 120, 22, 'Homeostasis Theory', 8),
        # 3) ___
        (148, 614, 110, 22, 'Evolution Theory', 9),
    ], notes='Page 12: Label exercise + 3 basic theories of biology.')

    # ═══════════════════════════════════════════════
    # PAGE 13: Experimental Design — Variables
    # ═══════════════════════════════════════════════
    build_page_slide(prs, 13, pages_dir / 'page-13.png', [
        # ___: a scientific investigation/procedure
        (62, 118, 100, 26, 'Experiment', 11),
        # ___: a changeable element
        (62, 146, 80, 24, 'Variable', 10),
        # ___ main types of variables
        (418, 170, 40, 22, 'two', 10),
        # 1) ___ Variable: Variable ___ by researcher
        (92, 224, 120, 22, 'Independent', 9),
        (422, 224, 100, 22, 'manipulated', 9),
        # 2) ___ Variable: Variable ___ by researcher
        (92, 258, 110, 22, 'Dependent', 9),
        (398, 258, 80, 22, 'measured', 9),
        # Independent Variable (___ of H2O)
        (338, 472, 65, 18, 'Amount', 8),
        # IV: style/model of paper airplane
        (170, 586, 200, 22, 'Style/model of paper airplane', 8),
        # DV: distance traveled
        (170, 610, 200, 22, 'Distance traveled', 8),
    ], notes='Page 13: Variables — independent (manipulated) vs dependent (measured).')

    # ═══════════════════════════════════════════════
    # PAGE 14: Controls & False Results
    # ═══════════════════════════════════════════════
    build_page_slide(prs, 14, pages_dir / 'page-14.png', [
        # contain ___ groups
        (328, 180, 80, 24, 'control', 10),
        # False ___: falsely indicate the ___ of a result
        (100, 210, 80, 24, 'Positives', 10),
        (518, 210, 80, 24, 'presence', 10),
        # False ___: falsely indicate the ___ of a result
        (100, 268, 80, 24, 'Negatives', 10),
        (518, 268, 80, 24, 'absence', 10),
        # ___ main types of controls
        (62, 336, 35, 24, 'Two', 10),
        # only differ in the ___ factor being tested
        (548, 362, 50, 22, 'one', 10),
        # 1) ___ Control: no response expected
        (92, 400, 80, 22, 'Negative', 10),
        (350, 400, 35, 22, 'no', 9),
        # 2) ___ Control: response expected
        (92, 432, 80, 22, 'Positive', 10),
    ], notes='Page 14: False positives/negatives, negative and positive controls.')

    # ═══════════════════════════════════════════════
    # PAGE 15: Controls practice (no blanks — just questions)
    # ═══════════════════════════════════════════════
    build_page_slide(prs, 15, pages_dir / 'page-15.png', [],
        notes='Page 15: Matching controls exercise + nitrogen experiment practice.')

"""


def main():
    output_path = Path(__file__).parent.parent / 'Campbell-Biology-Ch1-Slides.pptx'

    print(f"\n{'='*60}")
    print(f"  BUILDING SLIDE DECK v2: PDF Pages + Click Reveals")
    print(f"{'='*60}\n")

    prs = build_ch1()
    prs.save(str(output_path))

    print(f"  Slides: {len(prs.slides)}")
    print(f"  Output: {output_path}")
    print(f"  File size: {output_path.stat().st_size // 1024} KB")
    print(f"\n{'='*60}\n")


if __name__ == '__main__':
    main()
